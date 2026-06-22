from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "OmniNova_Claw_Investor_Pitch_CN.pptx"
LOGO = ROOT / "apps/omninova-tauri/public/omninoval-logo.png"

PRIMARY = RGBColor(79, 99, 235)
SECONDARY = RGBColor(99, 102, 241)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
BG = RGBColor(245, 247, 251)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(16, 185, 129)


SLIDES = [
    {
        "title": "OmniNova Claw",
        "subtitle": "本地优先的企业 AI 自动化工作台",
        "bullets": [
            "让 AI 从“会聊天”升级为“可执行、可审计、可私有化部署”的生产系统",
            "聚焦场景：客服自动化、运营自动化、后台流程执行",
            "关键词：本地优先 / 多渠道 / 技能系统 / 浏览器与桌面执行 / 安全治理",
        ],
    },
    {
        "title": "问题：企业 AI 还没有真正进入工作流",
        "bullets": [
            "大多数企业对 AI 的使用仍停留在问答、写稿、总结层面",
            "多渠道客服、后台录入、流程执行仍依赖人工切换多个系统",
            "企业最担心的是数据安全、权限失控、不可审计与落地复杂",
            "市场缺少真正能进入业务流程的 AI 执行系统",
        ],
        "highlight": "企业缺的不是另一个聊天工具，而是一个能完成任务的 AI 工作台。",
    },
    {
        "title": "解决方案：把模型、渠道、技能和执行整合成一套系统",
        "bullets": [
            "多模型接入：OpenAI、Anthropic、Gemini、DeepSeek、Ollama",
            "多渠道接入：Slack、Discord、Telegram、Webhook 等消息入口",
            "记忆与知识：工作记忆、情景记忆、技能/知识记忆",
            "工具执行：文件、命令、网页、浏览器、桌面控制",
            "安全治理：审批、权限策略、E-stop、执行审计",
        ],
    },
    {
        "title": "切入场景：先做客服与运营自动化",
        "bullets": [
            "场景一：多渠道客服自动化，统一接收消息、生成回复、人工复核、发送留痕",
            "场景二：运营自动化，自动执行后台录入、数据采集、报表生成、SOP 跟进",
            "场景三：内部知识与流程助手，将 FAQ、制度、SOP 技能化",
            "共同特征：高频、刚需、ROI 清晰、适合标准化模板复制",
        ],
        "highlight": "先切高频刚需场景，再扩展到更广泛的企业流程自动化。",
    },
    {
        "title": "产品形态：桌面端 + 本地网关 + CLI",
        "bullets": [
            "桌面控制中心：面向业务、运营与管理员的可视化配置入口",
            "本地网关：统一编排模型、技能、渠道与工具的执行中枢",
            "CLI：面向开发者与企业 IT 的自动化和部署入口",
            "技能系统：把企业 SOP、FAQ 与知识固化成可复用资产",
        ],
    },
    {
        "title": "为什么是现在：AI 执行层正处在窗口期",
        "bullets": [
            "大模型能力已足以支撑中等复杂度业务执行",
            "企业采购逻辑从“试用 AI”转向“验证 ROI”",
            "本地优先、私有化、安全治理正在成为企业落地前提",
            "聊天助手与传统 RPA 之间存在明显的执行层空白",
        ],
    },
    {
        "title": "差异化：不是聊天工具，也不是传统 RPA",
        "bullets": [
            "相比聊天工具：不仅能回答，还能执行",
            "相比传统 RPA：不仅会操作，还能理解上下文",
            "相比云端 Agent：更强调本地部署、私有化与企业可控性",
            "核心壁垒：本地优先架构、多模型多渠道编排、技能沉淀、安全控制",
        ],
        "highlight": "我们做的是企业 AI 的执行层，而不是内容生成壳子。",
    },
    {
        "title": "商业模式：先试点收费，再模板复制，再走企业订阅",
        "bullets": [
            "标准版：按团队、工作区或部署实例订阅",
            "企业版：私有化部署费 + 年服务费",
            "增值收入：行业技能包、定制集成、实施与培训",
            "收入路径：试点项目验证价值 → 模板化复制 → 订阅化扩张",
        ],
    },
    {
        "title": "当前进展与未来 12 个月目标",
        "bullets": [
            "已完成：桌面端、Rust 核心运行时、多模型支持、技能系统、本地网关、CLI",
            "目标一：落地 3–5 家试点客户",
            "目标二：打磨客服与运营两套标准模板",
            "目标三：跑通部署、续费与 ROI 指标",
            "目标四：形成私有化交付方法论",
        ],
    },
    {
        "title": "融资与用途：从可运行产品走向可复制业务",
        "bullets": [
            "本轮目标：验证 PMF 与场景复制能力",
            "资金用途：强化稳定性、安全审计、模板能力和客户成功体系",
            "市场目标：沉淀标杆客户案例，建立销售与实施闭环",
            "终局定位：成为企业 AI 进入真实工作流的执行入口",
        ],
        "highlight": "我们不是在造一个更会聊天的 AI，而是在造企业真正敢用的 AI 工作系统。",
    },
]


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_logo(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.7), Inches(0.45), height=Inches(0.7))


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(10.3), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(10.0), Inches(0.45))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = MUTED


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.2), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)
        p.bullet = True


def add_highlight(slide, text):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.1), Inches(11.4), Inches(0.85)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SECONDARY
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = SECONDARY
    p.alignment = PP_ALIGN.CENTER


def add_footer(slide, index, total):
    left = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(5.5), Inches(0.25))
    p = left.text_frame.paragraphs[0]
    p.text = "OmniNova Claw · Investor Pitch"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED

    right = slide.shapes.add_textbox(Inches(11.75), Inches(7.03), Inches(0.7), Inches(0.25))
    p2 = right.text_frame.paragraphs[0]
    p2.text = f"{index}/{total}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.RIGHT


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = len(SLIDES)

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_logo(slide)
        add_title(slide, data["title"], data.get("subtitle"))
        add_bullets(slide, data["bullets"])
        if data.get("highlight"):
            add_highlight(slide, data["highlight"])
        if idx == 1:
            metric = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.3), Inches(3.1), Inches(0.72)
            )
            metric.fill.solid()
            metric.fill.fore_color.rgb = GREEN
            metric.line.fill.background()
            tf = metric.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "先切客服与运营自动化"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        add_footer(slide, idx, total)

    prs.save(OUT)


if __name__ == "__main__":
    build()
